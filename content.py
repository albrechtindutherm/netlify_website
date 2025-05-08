from pptx import Presentation
import os

def pptx_to_markdown(input_pptx: str, output_md: str):
    prs = Presentation(input_pptx)
    with open(output_md, "w", encoding="utf-8") as md:
        md.write(f"# {os.path.splitext(os.path.basename(input_pptx))[0]}\n\n")
        for idx, slide in enumerate(prs.slides, 1):
            title = None
            for shape in slide.shapes:
                if shape.has_text_frame and 'title' in shape.name.lower():
                    title = shape.text.strip()
                    break
            header = title if title else f"Slide {idx}"
            md.write(f"## {header}\n\n")
            for shp in slide.shapes:
                if not shp.has_text_frame:
                    continue
                for paragraph in shp.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        md.write(f"- {text}\n")
            md.write("\n")

if __name__ == "__main__":
    pptx_to_markdown("presentation.pptx", "slides.md")
    print("Conversion complete: slides.md")
